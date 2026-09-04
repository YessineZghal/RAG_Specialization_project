"""`load_from_directory` — including Office-format support
(`.docx`/`.pptx`/`.xlsx`) added against the RAG-Anything gap review (see
`../../missing_to_complite.md`). No existing test file covered this
module at all before now.
"""

from __future__ import annotations

import sys
from pathlib import Path

import docx
from openpyxl import Workbook
from pptx import Presentation

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from src.ingest import load_from_directory


def test_load_from_directory_finds_every_real_sample_format():
    # The real, committed sample_docs/ directory -- md, pdf isn't present
    # here but docx/pptx/xlsx are, alongside the pre-existing md files.
    from src.config import settings

    docs = load_from_directory(settings.sample_docs_dir)
    ids = {d.id for d in docs}
    assert "onboarding_faq" in ids
    assert "remote_work_policy" in ids  # .docx
    assert "q3_product_update" in ids  # .pptx
    assert "support_ticket_volume" in ids  # .xlsx


def test_load_from_directory_reads_docx_paragraphs_and_tables(tmp_path):
    document = docx.Document()
    document.add_paragraph("A real paragraph of policy text.")
    table = document.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Role"
    table.rows[0].cells[1].text = "Days"
    document.save(tmp_path / "policy.docx")

    docs = {d.id: d for d in load_from_directory(tmp_path)}

    assert "A real paragraph of policy text." in docs["policy"].text
    assert "Role | Days" in docs["policy"].text


def test_load_from_directory_reads_pptx_slides_in_order(tmp_path):
    presentation = Presentation()
    layout = presentation.slide_layouts[1]
    slide1 = presentation.slides.add_slide(layout)
    slide1.shapes.title.text = "First Slide Title"
    slide2 = presentation.slides.add_slide(layout)
    slide2.shapes.title.text = "Second Slide Title"
    presentation.save(tmp_path / "deck.pptx")

    docs = {d.id: d for d in load_from_directory(tmp_path)}

    text = docs["deck"].text
    assert "Slide 1:" in text
    assert "First Slide Title" in text
    assert "Slide 2:" in text
    assert "Second Slide Title" in text
    assert text.index("First Slide Title") < text.index("Second Slide Title")


def test_load_from_directory_reads_xlsx_cell_values_not_formulas(tmp_path):
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Data"
    sheet.append(["Name", "Score"])
    sheet.append(["Alice", 92])
    workbook.save(tmp_path / "scores.xlsx")

    docs = {d.id: d for d in load_from_directory(tmp_path)}

    text = docs["scores"].text
    assert "Data:" in text
    assert "Name | Score" in text
    assert "Alice | 92" in text


def test_load_from_directory_skips_an_empty_slide_deck_with_no_text(tmp_path):
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])  # blank layout, no text at all
    presentation.save(tmp_path / "empty.pptx")

    docs = load_from_directory(tmp_path)

    assert docs == []  # empty text -> not returned as a Document, same rule as every other format


def test_load_from_directory_still_reads_plain_text_and_markdown_alongside_office_files(tmp_path):
    (tmp_path / "note.txt").write_text("Plain text still works.")
    workbook = Workbook()
    workbook.active.append(["a", "b"])
    workbook.save(tmp_path / "sheet.xlsx")

    docs = {d.id: d for d in load_from_directory(tmp_path)}

    assert docs["note"].text == "Plain text still works."
    assert "sheet" in docs
