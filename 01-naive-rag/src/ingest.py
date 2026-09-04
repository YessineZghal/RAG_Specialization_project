"""Load raw documents into the pipeline's common `Document` shape.

Two sources are supported:

1. `load_from_hf_dataset` — an **open-source** Hugging Face dataset,
   downloaded and cached by the `datasets` library the first time it runs.
   Default: `rag-datasets/rag-mini-wikipedia` (text-corpus/passages config),
   a small, public-domain-derived Wikipedia passage set built specifically
   for RAG tutorials. Nothing is bundled in this repo — see README.md#dataset.

2. `load_from_directory` — plain `.txt`/`.md`, `.pdf`, and Office files
   (`.docx`, `.pptx`, `.xlsx`) from a local folder (e.g.
   `data/sample_docs/`), for the "local PDF chatbot" mini project.
   Office support (see `../../missing_to_complite.md`'s RAG-Anything gap
   review) is deliberately pure-Python (`python-docx`/`python-pptx`/
   `openpyxl`), not a LibreOffice dependency — this repo's own
   `uv sync` is still the only setup step required.

Both return a flat `list[Document]` so the rest of the pipeline never has
to know where a document came from.
"""

from __future__ import annotations

import logging
from pathlib import Path

from pypdf import PdfReader

from .config import settings
from .schema import Document

logger = logging.getLogger(__name__)

TEXT_SUFFIXES = {".txt", ".md"}


def load_from_hf_dataset(
    dataset_name: str | None = None,
    config: str | None = None,
    split: str | None = None,
    limit: int | None = None,
) -> list[Document]:
    """Download (and cache) an open-source HF dataset and adapt it to `Document`.

    This is the only place in Level 1 that touches the network. It is never
    called at import time — only when a script or notebook explicitly runs
    it, and results are cached locally by `datasets` (in `~/.cache/huggingface`)
    so subsequent runs are offline.
    """
    from datasets import load_dataset  # imported lazily: heavy, network-capable

    dataset_name = dataset_name or settings.hf_dataset_name
    config = config or settings.hf_dataset_config
    split = split or settings.hf_dataset_split

    logger.info("Loading %s (%s/%s) from Hugging Face...", dataset_name, config, split)
    ds = load_dataset(dataset_name, config, split=split)
    if limit is not None:
        ds = ds.select(range(min(limit, len(ds))))

    documents: list[Document] = []
    for row in ds:
        # rag-mini-wikipedia's text-corpus config exposes `id` and `passage`.
        # Fall back gracefully if a different dataset is swapped in.
        doc_id = str(row.get("id", len(documents)))
        text = row.get("passage") or row.get("text") or row.get("content") or ""
        if not text.strip():
            continue
        documents.append(
            Document(
                id=f"hf-{doc_id}",
                text=text.strip(),
                metadata={"source": dataset_name, "split": split},
            )
        )

    logger.info("Loaded %d documents from %s.", len(documents), dataset_name)
    return documents


def load_from_directory(directory: Path | str | None = None) -> list[Document]:
    """Load local `.txt`, `.md`, and `.pdf` files as `Document`s.

    Used by the "local PDF chatbot" mini project and by the offline test
    suite (against the small hand-written files in `data/sample_docs/`).
    """
    directory = Path(directory) if directory is not None else settings.sample_docs_dir
    documents: list[Document] = []

    for path in sorted(directory.rglob("*")):
        if not path.is_file():
            continue

        suffix = path.suffix.lower()
        if suffix in TEXT_SUFFIXES:
            text = path.read_text(encoding="utf-8")
        elif suffix == ".pdf":
            text = _read_pdf(path)
        elif suffix == ".docx":
            text = _read_docx(path)
        elif suffix == ".pptx":
            text = _read_pptx(path)
        elif suffix == ".xlsx":
            text = _read_xlsx(path)
        else:
            continue

        if not text.strip():
            continue

        documents.append(
            Document(
                id=path.stem,
                text=text.strip(),
                metadata={"source": str(path.relative_to(directory.parent))},
            )
        )

    logger.info("Loaded %d local documents from %s.", len(documents), directory)
    return documents


def _read_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    return "\n\n".join(page.extract_text() or "" for page in reader.pages)


def _read_docx(path: Path) -> str:
    """Every paragraph's text, plus every table cell's text -- python-docx
    exposes both directly as structured objects (no `pdfplumber`-style
    mangling to work around, unlike this repo's PDF tables in Level 3)."""
    import docx

    document = docx.Document(str(path))
    parts = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip(" |"):
                parts.append(row_text)
    return "\n\n".join(parts)


def _read_pptx(path: Path) -> str:
    """Every text frame on every slide, in slide order -- speaker notes
    are deliberately excluded (they're presenter-facing, not part of what
    an audience actually sees, and would be a different, notes-specific
    corpus if this level ever wanted one)."""
    from pptx import Presentation

    presentation = Presentation(str(path))
    parts = []
    for slide_num, slide in enumerate(presentation.slides, start=1):
        slide_parts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        if slide_parts:
            parts.append(f"Slide {slide_num}: " + " | ".join(slide_parts))
    return "\n\n".join(parts)


def _read_xlsx(path: Path) -> str:
    """Every populated cell, row by row, per sheet -- `data_only=True` so
    a formula cell yields its last-computed value (what a reader actually
    sees), not the formula text itself."""
    import openpyxl

    workbook = openpyxl.load_workbook(str(path), data_only=True)
    parts = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows():
            cell_values = [str(cell.value) for cell in row if cell.value is not None]
            if cell_values:
                parts.append(f"{sheet.title}: " + " | ".join(cell_values))
    return "\n\n".join(parts)
